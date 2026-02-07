"""Document text extraction helpers for supported study-material formats."""
from __future__ import annotations

import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List

import olefile
import pymupdf
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".doc"}


@dataclass
class ExtractedDocument:
    text: str
    page_count: int
    extension: str


def infer_extension(filename: str) -> str:
    return Path(filename).suffix.lower()


def _normalize_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def _extract_pdf_text(file_path: str) -> ExtractedDocument:
    doc = pymupdf.open(file_path)
    try:
        text_parts: List[str] = []
        for index in range(doc.page_count):
            page_text = doc[index].get_text().strip()
            if page_text:
                text_parts.append(f"[Page {index + 1}]\n{page_text}")
        text = _normalize_text("\n\n".join(text_parts))
        return ExtractedDocument(
            text=text,
            page_count=max(doc.page_count, 1),
            extension=".pdf",
        )
    finally:
        doc.close()


def _extract_docx_text(file_path: str) -> ExtractedDocument:
    document = Document(file_path)
    text_parts: List[str] = []

    for paragraph in document.paragraphs:
        content = paragraph.text.strip()
        if content:
            text_parts.append(content)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                text_parts.append(" | ".join(cells))

    text = _normalize_text("\n\n".join(text_parts))
    estimated_pages = max(1, len(text) // 2400) if text else 1
    return ExtractedDocument(text=text, page_count=estimated_pages, extension=".docx")


def _extract_pptx_text(file_path: str) -> ExtractedDocument:
    presentation = Presentation(file_path)
    slide_text_blocks: List[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        shape_text: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                shape_text.append(text.strip())
        if shape_text:
            slide_text_blocks.append(f"[Slide {slide_index}]\n" + "\n".join(shape_text))

    text = _normalize_text("\n\n".join(slide_text_blocks))
    return ExtractedDocument(
        text=text,
        page_count=max(len(presentation.slides), 1),
        extension=".pptx",
    )


def _extract_doc_text(file_path: str) -> ExtractedDocument:
    """Best-effort extraction for legacy .doc files (OLE2 binary format)."""
    if not olefile.isOleFile(file_path):
        raise ValueError("Invalid .doc file: OLE header not found")

    extracted_fragments: List[str] = []
    with olefile.OleFileIO(file_path) as ole:
        for stream_name in ("WordDocument", "1Table", "0Table"):
            if not ole.exists(stream_name):
                continue
            raw = ole.openstream(stream_name).read()
            decoded = raw.decode("latin-1", errors="ignore")
            # Keep readable stretches and drop binary noise.
            readable = re.findall(r"[A-Za-z][A-Za-z0-9 ,.;:()\-_/]{20,}", decoded)
            extracted_fragments.extend(readable)

    text = _normalize_text("\n".join(extracted_fragments))
    if len(text) < 120:
        raise ValueError(
            "Could not reliably extract readable text from .doc. "
            "Please convert to .docx or .pdf for best results."
        )
    estimated_pages = max(1, len(text) // 2400)
    return ExtractedDocument(text=text, page_count=estimated_pages, extension=".doc")


def extract_document_text(file_path: str, filename: str) -> ExtractedDocument:
    extension = infer_extension(filename)
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file extension '{extension}'. "
            "Supported formats: .pdf, .docx, .pptx, .doc"
        )

    if extension == ".pdf":
        return _extract_pdf_text(file_path)
    if extension == ".docx":
        return _extract_docx_text(file_path)
    if extension == ".pptx":
        return _extract_pptx_text(file_path)
    return _extract_doc_text(file_path)

